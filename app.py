from flask import Flask, render_template, request, send_file
from pptx import Presentation
from fpdf import FPDF
import re
from groq import Groq
from pptx.util import Pt
import os
'''from dotenv import load_dotenv

load_dotenv()'''
app = Flask(__name__)

client = Groq(api_key=os.getenv("GROQ_API_KEY"))


generated_text = ""

# ---------------- PPT ----------------

# def create_ppt(content):

#     prs = Presentation()

#     # Strong splitter
#     slides = content.split("###SLIDE###")

#     for slide in slides:
#         slide = slide.strip()

#         if len(slide) < 10:
#             continue

#         lines = slide.split("\n")

#         title = lines[0]

#         s = prs.slides.add_slide(prs.slide_layouts[1])
#         s.shapes.title.text = title

#         tf = s.placeholders[1].text_frame
#         tf.clear()

#         for line in lines[1:]:
#             line = line.strip()
#             if line.startswith("-"):
#                 tf.add_paragraph().text = line.replace("-","").strip()

#     prs.save("output.pptx")
#     return "output.pptx"

def create_ppt(content):
    prs = Presentation()

    # Split slides by marker
    slides = content.split("###SLIDE###")

    for slide in slides:
        slide = slide.strip()
        if len(slide) < 5:  # skip very short slides
            continue

        lines = slide.split("\n")
        if not lines:
            continue

        title = lines[0].strip()
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = title

        tf = s.placeholders[1].text_frame
        tf.clear()

        # Add each line as paragraph, handle dashes/bullets
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
            # Remove common bullet characters if present
            if line[0] in ["-", "•"]:
                line = line[1:].strip()
            tf.add_paragraph().text = line

    prs.save("output.pptx")
    return "output.pptx"



# ---------------- PDF ----------------

# def create_pdf(content):

#     pdf = FPDF()
#     pdf.add_page()
#     pdf.set_font("Arial", size=12)

#     slides = re.split(r"###SLIDE###", content)

#     for slide in slides:
#         slide = slide.strip()
#         if not slide:
#             continue

#         lines = slide.split("\n")

#         pdf.set_font("Arial","B",14)
#         pdf.multi_cell(0,10,lines[0])

#         pdf.set_font("Arial", size=12)

#         for b in lines[1:]:
#             if b.strip():
#                 pdf.multi_cell(0,8,b.replace("-",""))

#         pdf.add_page()

#     pdf.output("output.pdf")
#     return "output.pdf"
def create_pdf(content):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)

    slides = re.split(r"###SLIDE###", content)

    for slide in slides:
        slide = slide.strip()
        if not slide:
            continue

        lines = slide.split("\n")
        pdf.set_font("Arial","B",14)
        pdf.multi_cell(0,10,lines[0].strip())  # Title

        pdf.set_font("Arial", size=12)
        for line in lines[1:]:
            line = line.strip()
            if not line:
                continue
            # Remove bullet characters
            if line[0] in ["-", "•"]:
                line = line[1:].strip()
            pdf.multi_cell(0,8,line)

        pdf.add_page()

    pdf.output("output.pdf")
    return "output.pdf"



# ---------------- ROUTES ----------------

@app.route("/", methods=["GET","POST"])
def home():

    global generated_text

    if request.method=="POST":

        topic = request.form["topic"]
        slides = request.form["slides"]
        format = request.form["format"]

        prompt = f"""
Create study PowerPoint on {topic}

Generate exactly {slides} slides.

STRICT FORMAT:

###SLIDE###
Introduction
- Full sentence.
- Proper grammar.
- Explain concepts.
- Academic tone.

###SLIDE###
Next Title
- Full sentence.
- Full sentence.
- Full sentence.
- Full sentence.

Rules:
• Always start each slide with ###SLIDE###
• First line is title
• Only bullets after
• 4 bullets per slide
• Last slide conclusion
"""


        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[{"role":"user","content":prompt}]
        )

        generated_text = response.choices[0].message.content

        return render_template("index.html", text=generated_text, format=format)

    return render_template("index.html")

@app.route("/download")
def download():

    global generated_text
    format = request.args.get("format")

    if format=="pdf":
        file = create_pdf(generated_text)
    else:
        file = create_ppt(generated_text)

    return send_file(file, as_attachment=True)

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 10000)))

